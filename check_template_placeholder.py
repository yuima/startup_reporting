from pptx import Presentation

def inspect_placeholders(template_path):
    prs = Presentation(template_path)
    slide = prs.slides[0]  # テンプレートの1枚目のスライドを確認
    for shape in slide.shapes:
        if shape.is_placeholder:
            print(f"インデックス: {shape.placeholder_format.idx}, 名前: {shape.name}")

# 実行例
inspect_placeholders("template_startup_report.pptx")
