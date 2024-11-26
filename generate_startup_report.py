import os
from pptx import Presentation
import tkinter as tk
from tkinter import filedialog


def create_ppt_with_template_and_placeholders(data, template_path, output_folder):
    # テンプレートを読み込む
    prs = Presentation(template_path)

    # 最初のスライドを取得（テンプレートの1枚目のスライドを利用）
    slide = prs.slides[0]

    # プレースホルダーのインデックスに基づいてデータを挿入
    for shape in slide.shapes:
        if shape.is_placeholder:
            idx = shape.placeholder_format.idx
            if idx == 0:  # タイトルプレースホルダー
                shape.text = f"スタートアップレポート ({data['企業名']})"
            elif idx == 13:  # 本社
                shape.text = f"本社: {data['本社']}"
            elif idx == 14:  # 設立年
                shape.text = f"設立年: {data['設立年']}"
            elif idx == 15:  # URL
                shape.text = f"URL: {data['URL']}"
            elif idx == 16:  # ステージ
                shape.text = f"ステージ: {data['ステージ']}"
            elif idx == 17:  # ボードメンバ名
                shape.text = f"ボードメンバ名: {data['ボードメンバ名']}"
            elif idx == 18:  # 主要投資家
                shape.text = f"主要投資家: {data['主要投資家']}"
            elif idx == 19:  # 企業概要
                shape.text = f"企業概要: {data['企業概要']}"
            elif idx == 20:  # 競合
                shape.text = f"競合: {data['競合']}"
            elif idx == 21:  # 強みと機会
                shape.text = f"強みと機会: {data['強みと機会']}"
            elif idx == 21:  # 弱みと脅威
                shape.text = f"弱みと脅威: {data['弱みと脅威']}"
            elif idx == 22:  # 主要顧客やビジネスモデル
                shape.text = f"主要顧客やビジネスモデル: {data['主要顧客やビジネスモデル']}"

    # 企業名を元に出力ファイル名を生成
    company_name = data.get('企業名', '未定').replace(" ", "_").replace("/", "_")
    output_path = os.path.join(output_folder, f"{company_name}_スタートアップレポート.pptx")

    # 保存
    prs.save(output_path)
    print(f"PPTファイルを作成しました: {output_path}")


def read_text_file(file_path):
    """
    テキストファイルからデータを読み込み、辞書として返す
    :param file_path: str, テキストファイルのパス
    :return: dict, テキストデータを辞書形式で返す
    """
    data = {}
    with open(file_path, 'r', encoding='utf-8') as file:
        lines = file.readlines()
        for line in lines:
            if ":" in line:
                key, value = line.split(":", 1)
                data[key.strip()] = value.strip()
    return data


def select_file():
    """ファイル選択ダイアログを表示"""
    tk.Tk().withdraw()  # GUIウィンドウを非表示
    file_path = filedialog.askopenfilename(
        title="テキストファイルを選択してください",
        filetypes=[("Text Files", "*.txt"), ("All Files", "*.*")]
    )
    return file_path


def select_file(title, filetypes):
    """ファイル選択ダイアログを表示"""
    tk.Tk().withdraw()  # GUIウィンドウを非表示
    file_path = filedialog.askopenfilename(title=title, filetypes=filetypes)
    return file_path

def main():
    print("### スタートアップレポート生成ツール ###")
    
    # テキストファイルの選択
    print("テキストファイルを選択してください...")
    input_file = select_file("テキストファイルを選択してください", [("Text Files", "*.txt"), ("All Files", "*.*")])
    if not input_file:
        print("エラー: ファイルが選択されませんでした。")
        input("エンターキーを押して終了します...")
        return

    # テンプレートファイルの選択
    print("テンプレートファイルを選択してください...")
    template_path = select_file("テンプレートファイルを選択してください", [("PowerPoint Templates", "*.pptx"), ("All Files", "*.*")])
    if not template_path:
        print("エラー: テンプレートファイルが選択されませんでした。")
        input("エンターキーを押して終了します...")
        return

    # 出力フォルダを確認・作成
    output_folder = "output"
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # データを読み込んでPPTを生成
    data = read_text_file(input_file)
    create_ppt_with_template_and_placeholders(data, template_path, output_folder)

    print("完了しました。エンターキーを押して終了します...")
    input()


if __name__ == "__main__":
    main()

# # 実行
# input_file = "input/sample.txt"  # ユーザー入力のテキストファイル
# template_path = "template_startup_report.pptx"  # テンプレートファイルのパス
# output_path = "output/"  # 出力ファイル名
# data = read_text_file(input_file)
# create_ppt_with_template_and_placeholders(data, template_path, output_path)
